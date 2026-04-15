"""
PPT estilo SLIDE_RESUMEN.html - cards redondeadas, sin preguntas
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colores ────────────────────────────────────────────────────────
NAVY     = RGBColor(0x00, 0x1e, 0x60)
BLUE     = RGBColor(0x00, 0x53, 0xe2)
SPARK    = RGBColor(0xFF, 0xC2, 0x20)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG = RGBColor(0xF0, 0xF4, 0xFB)
CARD_BG  = RGBColor(0xF8, 0xF9, 0xFB)
GREEN    = RGBColor(0x2a, 0x87, 0x03)
RED      = RGBColor(0xEA, 0x11, 0x00)
CYAN     = RGBColor(0x00, 0x77, 0xAA)
AMBER    = RGBColor(0x99, 0x52, 0x13)
GREEN_BG = RGBColor(0xEA, 0xF6, 0xEC)
YELL_BG  = RGBColor(0xFF, 0xF9, 0xE6)
TEXT     = RGBColor(0x44, 0x44, 0x55)
LINE     = RGBColor(0xE0, 0xE6, 0xF5)

ROUNDED = 5   # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
RECT    = 1   # MSO_AUTO_SHAPE_TYPE.RECTANGLE


def add_shape(slide, kind, l, t, w, h, fill=None, line_color=None,
              line_w=Pt(0), radius=0.05):
    sp = slide.shapes.add_shape(
        kind, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = line_w
    else:
        sp.line.fill.background()
    if kind == ROUNDED:
        sp.adjustments[0] = radius
    return sp


def txb(slide, text, l, t, w, h,
        size=11, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color or NAVY
    r.font.name = "Calibri"
    return tb


# ── Presentacion ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Fondo
add_shape(slide, RECT, 0, 0, 13.33, 7.5, SLIDE_BG)

# ── HEADER REDONDEADO (solo esquinas abajo) ────────────────────────
# Usamos rect normal para el header (va de borde a borde)
add_shape(slide, RECT, 0, 0, 13.33, 1.22, NAVY)
# Franja spark debajo del header
add_shape(slide, RECT, 0, 1.18, 13.33, 0.07, SPARK)

# Logo spark (estrella)
txb(slide, "*", 0.28, 0.1, 0.55, 0.75,
    size=36, bold=True, color=SPARK)

# Titulos
txb(slide, "Cuestionario PyME  Walmart Chile 2026",
    0.88, 0.08, 9.0, 0.58,
    size=23, bold=True, color=WHITE)
txb(slide, "Resumen de bloques  |  25 preguntas  |  5 areas clave  |  menos de 4 minutos",
    0.88, 0.65, 9.5, 0.38,
    size=10, color=RGBColor(0xAA, 0xBB, 0xFF))

# Pill CONFIDENCIAL (redondeado)
add_shape(slide, ROUNDED, 10.75, 0.36, 2.32, 0.38,
          fill=SPARK, radius=0.5)
txb(slide, "CONFIDENCIAL  |  USO INTERNO",
    10.75, 0.36, 2.32, 0.38,
    size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── STATS ROW (4 cards redondeadas) ───────────────────────────────
stats = [
    ("25",    "TOTAL PREGUNTAS",  NAVY,     WHITE),
    ("46",    "PUNTAJE MAXIMO",   SPARK,    NAVY),
    (">= 40", "PYME MADURA (pts)",  GREEN_BG, GREEN),
    ("< 40",  "PYME INTERMEDIA (pts)", YELL_BG, RGBColor(0x6b, 0x3a, 0x00)),
]
sw = 2.9
sgap = 0.22
sx = (13.33 - sw * 4 - sgap * 3) / 2

for i, (num, lbl, bg_c, fg_c) in enumerate(stats):
    x = sx + i * (sw + sgap)
    add_shape(slide, ROUNDED, x, 1.33, sw, 0.78,
              fill=bg_c, radius=0.08)
    txb(slide, num, x, 1.34, sw, 0.44,
        size=22, bold=True, color=fg_c, align=PP_ALIGN.CENTER)
    txb(slide, lbl, x, 1.74, sw, 0.28,
        size=8, bold=True, color=fg_c, align=PP_ALIGN.CENTER)

# ── CARDS GRID 3x2 ────────────────────────────────────────────────
cw = 4.16   # card width
ch = 2.12   # card height
cgx = 0.16  # gap horizontal
cgy = 0.14  # gap vertical
cs_x = 0.23
cs_y = 2.26

cards = [
    (
        "Bloque 1  |  Identificacion",
        "Datos de la empresa, categoria de producto, "
        "antiguedad como proveedor e ingreso al Programa Potencia PyME.",
        "Max: 4 pts",
        BLUE,
        RGBColor(0xEE, 0xF3, 0xFD),
    ),
    (
        "Bloque 2  |  Supply Chain & CD",
        "Fill Rate, In-Stock, requisitos de entrega en el CD, "
        "agendamiento de citas, inventario en sala y quiebres de stock.",
        "Max: 12 pts",
        CYAN,
        RGBColor(0xE8, 0xF5, 0xFB),
    ),
    (
        "Bloque 3  |  Retail Link & Phronesis",
        "Acceso a Retail Link, frecuencia de analisis de ventas, "
        "uso de Phronesis, presencia por local y alzas de costo.",
        "Max: 12 pts",
        AMBER,
        RGBColor(0xFB, 0xF5, 0xEA),
    ),
    (
        "Bloque 4  |  Gestion Comercial",
        "Promociones, costo neto vs limpio, representante de "
        "cuenta (KAM), Marca Propia y canales de contacto.",
        "Max: 10 pts",
        GREEN,
        RGBColor(0xEA, 0xF6, 0xEC),
    ),
    (
        "Bloque 5  |  Gestion Operacional",
        "Rechazos en CD, Fee de Ultima Milla, "
        "Fee de Reposicion y certificaciones de calidad (GFSI).",
        "Max: 8 pts",
        RED,
        RGBColor(0xFD, 0xEE, 0xEE),
    ),
    (
        "Escala de Puntaje  |  Total: 46 pts",
        "0 pts  No conoce / No lo hace\n"
        "1 pt    Conoce pero no gestiona activamente\n"
        "2 pts  Domina y gestiona con regularidad",
        "Umbral: 40 pts para ser Madura",
        NAVY,
        RGBColor(0xEE, 0xF2, 0xFF),
    ),
]

for i, (titulo, desc, pts_label, border_c, bg_c) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = cs_x + col * (cw + cgx)
    y = cs_y + row * (ch + cgy)

    # Card redondeada fondo
    add_shape(slide, ROUNDED, x, y, cw, ch,
              fill=bg_c, radius=0.06)

    # Header coloreado redondeado arriba
    add_shape(slide, ROUNDED, x, y, cw, 0.44,
              fill=border_c, radius=0.06)
    # Tapa la parte baja del header redondeado (para que sea recto abajo)
    add_shape(slide, RECT, x, y + 0.25, cw, 0.19,
              fill=border_c)

    # Titulo en blanco sobre header
    txb(slide, titulo,
        x + 0.18, y + 0.04, cw - 1.1, 0.36,
        size=10, bold=True, color=WHITE)

    # Badge puntos (esquina derecha del header)
    add_shape(slide, ROUNDED, x + cw - 0.88, y + 0.06, 0.82, 0.28,
              fill=WHITE, radius=0.5)
    txb(slide, pts_label,
        x + cw - 0.88, y + 0.05, 0.82, 0.28,
        size=7, bold=True, color=border_c, align=PP_ALIGN.CENTER)

    # Descripcion
    txb(slide, desc,
        x + 0.18, y + 0.52, cw - 0.28, ch - 0.65,
        size=10, color=TEXT)

# ── BARRA SEGMENTACION ────────────────────────────────────────────
bar_y = 6.67
txb(slide, "Segmentacion al completar el cuestionario:",
    0.23, bar_y - 0.23, 5.5, 0.2,
    size=9, bold=True, color=NAVY)

add_shape(slide, ROUNDED, 0.23, bar_y, 4.5, 0.38,
          fill=YELL_BG, line_color=SPARK, line_w=Pt(1.5), radius=0.4)
add_shape(slide, ROUNDED, 4.83, bar_y, 8.27, 0.38,
          fill=GREEN_BG, line_color=GREEN, line_w=Pt(1.5), radius=0.4)

txb(slide, "PyME INTERMEDIA   0 - 39 pts",
    0.23, bar_y + 0.01, 4.5, 0.35,
    size=9, bold=True,
    color=RGBColor(0x6b, 0x3a, 0x00), align=PP_ALIGN.CENTER)
txb(slide, "PyME MADURA   40 - 46 pts",
    4.83, bar_y + 0.01, 8.27, 0.35,
    size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ── FOOTER ────────────────────────────────────────────────────────
add_shape(slide, RECT, 0, 7.17, 13.33, 0.33, NAVY)
txb(slide, "Walmart Chile  |  Programa PyME 2026",
    0.3, 7.19, 5.0, 0.25, size=9, color=WHITE)
txb(slide, "pyme-survey.onrender.com/formulario",
    4.5, 7.19, 5.0, 0.25, size=9,
    color=RGBColor(0xAA, 0xBB, 0xFF), align=PP_ALIGN.CENTER)
txb(slide, "Documento de uso interno",
    8.5, 7.19, 4.5, 0.25, size=9,
    color=WHITE, align=PP_ALIGN.RIGHT)


OUT = r"C:\Users\c0b0w1w\Documents\puppy_workspace\pyme_survey\PyME_Resumen_v3.pptx"
prs.save(OUT)
print("OK:", OUT)